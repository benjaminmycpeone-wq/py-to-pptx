"""Sample script to test the converter. Upload this file to generate a PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Hello from Python!"
slide.placeholders[1].text = "Generated with python-pptx"

# Content slide
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "Key Points"
body = slide2.placeholders[1]
body.text = "This presentation was auto-generated"
body.text_frame.add_paragraph().text = "From a simple .py script"
body.text_frame.add_paragraph().text = "Using the PY to PPTX converter"

prs.save("demo.pptx")
print("Created demo.pptx")
