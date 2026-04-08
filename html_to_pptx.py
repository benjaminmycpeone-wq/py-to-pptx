"""Convert an HTML file into a PowerPoint presentation.

Strategy:
- Each <h1>/<h2> starts a new slide as the title.
- <p>, <li>, <pre>, <code> become body text on the current slide.
- <table> rows are rendered as a native PPTX table.
- If no headings exist, all content goes on a single slide titled after the filename.
"""

import os
import re
from bs4 import BeautifulSoup, NavigableString
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN


def html_to_pptx(html_path, output_dir=None):
    """Convert *html_path* to a .pptx file. Returns the output path."""
    with open(html_path, encoding="utf-8", errors="replace") as f:
        html = f.read()

    soup = BeautifulSoup(html, "html.parser")

    # Decide output path
    base = os.path.splitext(os.path.basename(html_path))[0]
    out_dir = output_dir or os.path.dirname(html_path)
    out_path = os.path.join(out_dir, f"{base}.pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = _parse_slides(soup, base)

    for slide_info in slides_data:
        _add_slide(prs, slide_info)

    prs.save(out_path)
    return out_path


# ---------------------------------------------------------------------------
# Parsing
# ---------------------------------------------------------------------------

def _parse_slides(soup, fallback_title):
    """Return a list of dicts: {title, items} where items are (type, data)."""
    slides = []
    current = None

    # Get the document title if available
    title_tag = soup.find("title")
    doc_title = title_tag.get_text(strip=True) if title_tag else fallback_title

    body = soup.find("body") or soup

    for el in body.children:
        if isinstance(el, NavigableString):
            text = el.strip()
            if text and current:
                current["items"].append(("text", text))
            continue

        tag = el.name
        if not tag:
            continue

        # Headings start a new slide
        if tag in ("h1", "h2", "h3"):
            current = {"title": el.get_text(strip=True), "items": []}
            slides.append(current)
            continue

        # Ensure we have at least one slide
        if current is None:
            current = {"title": doc_title, "items": []}
            slides.append(current)

        if tag in ("p", "div", "span", "blockquote", "address"):
            text = el.get_text(strip=True)
            if text:
                current["items"].append(("text", text))

        elif tag in ("ul", "ol"):
            items = []
            for li in el.find_all("li", recursive=False):
                items.append(li.get_text(strip=True))
            if items:
                current["items"].append(("list", items))

        elif tag in ("pre", "code"):
            current["items"].append(("code", el.get_text()))

        elif tag == "table":
            rows = _parse_table(el)
            if rows:
                current["items"].append(("table", rows))

        elif tag in ("h4", "h5", "h6"):
            text = el.get_text(strip=True)
            if text:
                current["items"].append(("subheading", text))

        elif tag == "img":
            alt = el.get("alt", "")
            if alt:
                current["items"].append(("text", f"[Image: {alt}]"))

        else:
            # Recurse into unknown containers (section, article, main, etc.)
            inner = _parse_slides(el, fallback_title)
            if inner:
                for s in inner:
                    if current and not current["items"] and not slides[-1]["items"]:
                        # Merge into empty current slide
                        current["title"] = s["title"] or current["title"]
                        current["items"] = s["items"]
                    else:
                        slides.append(s)
                        current = s

    # Fallback: empty presentation
    if not slides:
        slides.append({"title": doc_title, "items": [("text", "(empty document)")]})

    return slides


def _parse_table(table_el):
    rows = []
    for tr in table_el.find_all("tr"):
        cells = []
        for td in tr.find_all(["td", "th"]):
            cells.append(td.get_text(strip=True))
        if cells:
            rows.append(cells)
    return rows


# ---------------------------------------------------------------------------
# Slide creation
# ---------------------------------------------------------------------------

def _add_slide(prs, slide_info):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    title = slide_info["title"]
    items = slide_info["items"]

    # Title text box
    left = Inches(0.7)
    top = Inches(0.4)
    width = Inches(11.9)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _rgb(0x1E, 0x29, 0x3B)

    # Content area
    content_top = Inches(1.5)
    content_left = Inches(0.7)
    content_width = Inches(11.9)
    y = content_top

    for item_type, data in items:
        if y > Inches(6.8):
            break  # avoid overflow

        if item_type == "text":
            box = slide.shapes.add_textbox(content_left, y, content_width, Inches(0.6))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data
            p.font.size = Pt(16)
            p.font.color.rgb = _rgb(0x33, 0x41, 0x55)
            y += Inches(0.5)

        elif item_type == "subheading":
            box = slide.shapes.add_textbox(content_left, y, content_width, Inches(0.5))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = _rgb(0x1E, 0x29, 0x3B)
            y += Inches(0.55)

        elif item_type == "list":
            box = slide.shapes.add_textbox(content_left, y, content_width, Inches(len(data) * 0.35 + 0.2))
            tf = box.text_frame
            tf.word_wrap = True
            for i, li_text in enumerate(data):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"\u2022  {li_text}"
                p.font.size = Pt(15)
                p.font.color.rgb = _rgb(0x33, 0x41, 0x55)
                p.space_after = Pt(4)
            y += Inches(len(data) * 0.35 + 0.2)

        elif item_type == "code":
            box = slide.shapes.add_textbox(content_left, y, content_width, Inches(1.5))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data.strip()[:800]  # cap long code blocks
            p.font.size = Pt(11)
            p.font.name = "Consolas"
            p.font.color.rgb = _rgb(0x1E, 0x29, 0x3B)
            y += Inches(1.6)

        elif item_type == "table":
            rows_data = data
            n_rows = min(len(rows_data), 15)
            n_cols = max(len(r) for r in rows_data) if rows_data else 1
            col_w = min(content_width / n_cols, Inches(3))
            tbl_h = Inches(n_rows * 0.35)

            table_shape = slide.shapes.add_table(
                n_rows, n_cols, content_left, y, Emu(int(col_w * n_cols)), tbl_h
            )
            table = table_shape.table

            for r_idx in range(n_rows):
                for c_idx in range(min(n_cols, len(rows_data[r_idx]))):
                    cell = table.cell(r_idx, c_idx)
                    cell.text = rows_data[r_idx][c_idx]
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(11)

            y += tbl_h + Inches(0.3)


def _rgb(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python html_to_pptx.py <file.html>")
        sys.exit(1)
    out = html_to_pptx(sys.argv[1])
    print(f"Created {out}")
